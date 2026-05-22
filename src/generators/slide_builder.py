"""单页幻灯片生成器 - 基于python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import yaml
from pathlib import Path
import json


async def build_slide(
    blueprint_yaml: str,
    slide_index: int,
    modifications: str = "",
    template_path: str = None
) -> dict:
    """根据蓝图生成单页幻灯片。

    Args:
        blueprint_yaml: 完整蓝图YAML
        slide_index: 页码（从0开始）
        modifications: 修改意见
        template_path: 模板PPTX路径

    Returns:
        dict: {"pptx_path": str, "slide_index": int}
    """
    # 解析蓝图
    blueprint = yaml.safe_load(blueprint_yaml)

    if slide_index >= len(blueprint.get("slides", [])):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide_def = blueprint["slides"][slide_index]

    # 创建或使用模板
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # 设置为宽屏 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # 清空默认幻灯片
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # 添加幻灯片
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 根据类型填充内容
    slide_type = slide_def.get("type", "content")

    if slide_type == "title":
        _build_title_slide(slide, slide_def, prs)
    elif slide_type == "content":
        _build_content_slide(slide, slide_def, prs)
    elif slide_type == "chart":
        _build_chart_slide(slide, slide_def, prs)
    elif slide_type == "conclusion":
        _build_conclusion_slide(slide, slide_def, prs)
    else:
        _build_content_slide(slide, slide_def, prs)

    # 保存
    output_dir = Path("workspace/preview")
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / f"slide_{slide_index}.pptx"
    prs.save(str(pptx_path))

    return {
        "pptx_path": str(pptx_path),
        "slide_index": slide_index
    }


def _build_title_slide(slide, slide_def: dict, prs: Presentation):
    """构建标题页."""
    title = slide_def.get("title", "")
    subtitle = slide_def.get("subtitle", "")

    # 标题
    left, top, width, height = Inches(1), Inches(2), Inches(11), Inches(2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # 副标题
    left, top, width, height = Inches(1), Inches(4), Inches(11), Inches(1.5)
    sub_box = slide.shapes.add_textbox(left, top, width, height)
    sub_frame = sub_box.text_frame
    sub_frame.text = subtitle
    sub_para = sub_frame.paragraphs[0]
    sub_para.font.size = Pt(24)
    sub_para.alignment = PP_ALIGN.CENTER


def _build_content_slide(slide, slide_def: dict, prs: Presentation):
    """构建内容页."""
    title = slide_def.get("title", "")
    content = slide_def.get("content", "")

    # 标题
    left, top, width, height = Inches(0.5), Inches(0.3), Inches(12), Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # 内容
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    content_para = content_frame.paragraphs[0]
    content_para.text = content
    content_para.font.size = Pt(18)

    # 添加内容块
    content_blocks = slide_def.get("content_blocks", [])
    if content_blocks:
        content_frame.clear()
        for block in content_blocks:
            if isinstance(block, dict) and "text" in block:
                para = content_frame.add_paragraph()
                para.text = block["text"]
                para.font.size = Pt(18)


def _build_chart_slide(slide, slide_def: dict, prs: Presentation):
    """构建图表页."""
    title = slide_def.get("title", "")
    chart_desc = slide_def.get("chart", "")

    # 标题
    left, top, width, height = Inches(0.5), Inches(0.3), Inches(12), Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # 图表描述
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)
    desc_box = slide.shapes.add_textbox(left, top, width, height)
    desc_frame = desc_box.text_frame
    desc_frame.text = chart_desc
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(18)


def _build_conclusion_slide(slide, slide_def: dict, prs: Presentation):
    """构建总结页."""
    title = slide_def.get("title", "Conclusion")
    key_points = slide_def.get("key_points", [])
    future_work = slide_def.get("future_work", "")

    # 标题
    left, top, width, height = Inches(0.5), Inches(0.3), Inches(12), Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # 关键点
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(12), Inches(3.5)
    points_box = slide.shapes.add_textbox(left, top, width, height)
    points_frame = points_box.text_frame
    points_frame.word_wrap = True

    if key_points:
        for i, point in enumerate(key_points):
            if i == 0:
                para = points_frame.paragraphs[0]
            else:
                para = points_frame.add_paragraph()
            para.text = f"- {point}"
            para.font.size = Pt(18)

    # 未来工作
    if future_work:
        left, top, width, height = Inches(0.5), Inches(5), Inches(12), Inches(2)
        future_box = slide.shapes.add_textbox(left, top, width, height)
        future_frame = future_box.text_frame
        future_frame.text = f"Future Work: {future_work}"
        future_para = future_frame.paragraphs[0]
        future_para.font.size = Pt(16)
        future_para.font.italic = True
