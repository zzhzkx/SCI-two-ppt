"""PPTX到HTML转换器 - 实现预览与PPTX一致."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pathlib import Path


class PptxToHtmlConverter:
    """将PPTX文件转换为HTML预览，保持样式一致."""

    def __init__(self, pptx_path: str):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(pptx_path)

    def convert_slide_to_html(self, slide_index: int) -> str:
        """将单页幻灯片转换为HTML."""
        if slide_index >= len(self.prs.slides):
            return "<html><body>Slide not found</body></html>"

        slide = self.prs.slides[slide_index]

        # 获取幻灯片尺寸
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height

        # 提取背景
        background_html = self._extract_background(slide)

        # 提取所有形状
        shapes_html = self._extract_shapes(slide)

        # 生成完整HTML
        html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Slide {slide_index}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: Arial, sans-serif; background: #f5f5f5; padding: 20px; }}
        .slide-container {{
            width: {self._emu_to_px(slide_width)}px;
            height: {self._emu_to_px(slide_height)}px;
            margin: 0 auto;
            position: relative;
            overflow: hidden;
            {background_html}
        }}
        .shape {{
            position: absolute;
            word-wrap: break-word;
        }}
        .controls {{ margin-top: 20px; text-align: center; }}
        .btn {{ padding: 10px 20px; margin: 5px; cursor: pointer; border: none; border-radius: 4px; }}
        .btn-confirm {{ background: #27ae60; color: white; }}
        .btn-modify {{ background: #f39c12; color: white; }}
        .btn-skip {{ background: #95a5a6; color: white; }}
        .btn-redo {{ background: #e74c3c; color: white; }}
    </style>
</head>
<body>
    <div class="slide-container">
        {shapes_html}
    </div>
    <div class="controls">
        <button class="btn btn-confirm" onclick="sendAction('confirm')">Confirm</button>
        <button class="btn btn-modify" onclick="sendAction('modify')">Modify</button>
        <button class="btn btn-skip" onclick="sendAction('skip')">Skip</button>
        <button class="btn btn-redo" onclick="sendAction('redo')">Redo</button>
    </div>
    <script>
        function sendAction(action) {{
            fetch('/api/action', {{
                method: 'POST',
                headers: {{'Content-Type': 'application/json'}},
                body: JSON.stringify({{slide: {slide_index}, action: action}})
            }});
        }}
    </script>
</body>
</html>"""
        return html

    def _extract_background(self, slide) -> str:
        """提取背景样式."""
        background = slide.background
        fill = background.fill

        if fill.type is not None:
            if fill.type == 1:  # Solid
                color = fill.fore_color.rgb
                return f"background: #{color};"
            elif fill.type == 2:  # Gradient
                # 简化处理：使用第一个颜色
                try:
                    color = fill.gradient_stops[0].color.rgb
                    return f"background: #{color};"
                except:
                    return "background: #FFFFFF;"

        return "background: #FFFFFF;"

    def _extract_shapes(self, slide) -> str:
        """提取所有形状."""
        html_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                # 文本框
                left = self._emu_to_px(shape.left)
                top = self._emu_to_px(shape.top)
                width = self._emu_to_px(shape.width)
                height = self._emu_to_px(shape.height)

                text = shape.text_frame.text
                font_size = self._get_font_size(shape)
                font_color = self._get_font_color(shape)
                alignment = self._get_alignment(shape)
                bold = self._is_bold(shape)

                style = f"left:{left}px; top:{top}px; width:{width}px; height:{height}px;"
                style += f" font-size:{font_size}px; color:{font_color};"
                style += f" text-align:{alignment};"
                if bold:
                    style += " font-weight:bold;"

                html_parts.append(f'<div class="shape" style="{style}">{text}</div>')

        return "\n        ".join(html_parts)

    def _emu_to_px(self, emu: int) -> int:
        """EMU转像素 (1 inch = 96px, 1 inch = 914400 EMU)."""
        return int(emu * 96 / 914400)

    def _get_font_size(self, shape) -> int:
        """获取字体大小."""
        try:
            if shape.text_frame.paragraphs[0].runs:
                run = shape.text_frame.paragraphs[0].runs[0]
                if run.font.size:
                    return int(run.font.size.pt)
        except:
            pass
        return 18

    def _get_font_color(self, shape) -> str:
        """获取字体颜色."""
        try:
            if shape.text_frame.paragraphs[0].runs:
                run = shape.text_frame.paragraphs[0].runs[0]
                if run.font.color and run.font.color.rgb:
                    return f"#{run.font.color.rgb}"
        except:
            pass
        return "#000000"

    def _get_alignment(self, shape) -> str:
        """获取对齐方式."""
        try:
            para = shape.text_frame.paragraphs[0]
            if para.alignment == 1:  # CENTER
                return "center"
            elif para.alignment == 2:  # RIGHT
                return "right"
        except:
            pass
        return "left"

    def _is_bold(self, shape) -> bool:
        """检查是否加粗."""
        try:
            if shape.text_frame.paragraphs[0].runs:
                run = shape.text_frame.paragraphs[0].runs[0]
                return run.font.bold
        except:
            pass
        return False

    def save_html(self, slide_index: int, output_path: str):
        """保存HTML预览文件."""
        html = self.convert_slide_to_html(slide_index)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
