"""预览生成器 - 生成HTML/图片/PPT预览."""

from pathlib import Path
import json


class PreviewGenerator:
    """预览生成器，支持多种预览格式."""

    def __init__(self, output_dir: str = "workspace/preview"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_html_preview(self, slide_data: dict, slide_index: int) -> str:
        """生成HTML预览文件.

        Args:
            slide_data: 幻灯片数据
            slide_index: 幻灯片索引

        Returns:
            HTML文件路径
        """
        html_content = self._create_html_template(slide_data, slide_index)
        html_path = self.output_dir / f"slide_{slide_index}.html"

        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html_content)

        return str(html_path)

    def generate_image_preview(self, pptx_path: str, slide_index: int) -> str:
        """生成图片预览.

        Args:
            pptx_path: PPTX文件路径
            slide_index: 幻灯片索引

        Returns:
            图片文件路径
        """
        # 使用python-pptx渲染为图片
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(pptx_path)
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            # 这里需要使用PIL或其他库渲染为图片
            # 简化实现：返回占位路径
            image_path = self.output_dir / f"slide_{slide_index}.png"
            return str(image_path)

        return ""

    def _create_html_template(self, slide_data: dict, slide_index: int) -> str:
        """创建HTML模板."""
        title = slide_data.get("title", f"Slide {slide_index}")
        content = slide_data.get("content", "")
        notes = slide_data.get("notes", "")
        duration = slide_data.get("duration_seconds", 60)

        html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title} - Preview</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}

        body {{
            font-family: 'Segoe UI', Arial, sans-serif;
            background: #f5f5f5;
            padding: 20px;
        }}

        .container {{
            max-width: 1200px;
            margin: 0 auto;
        }}

        .header {{
            background: #003366;
            color: white;
            padding: 20px;
            border-radius: 8px 8px 0 0;
        }}

        .header h1 {{
            font-size: 24px;
            margin-bottom: 10px;
        }}

        .header .meta {{
            font-size: 14px;
            opacity: 0.8;
        }}

        .slide-preview {{
            background: white;
            border: 1px solid #ddd;
            padding: 40px;
            min-height: 500px;
        }}

        .slide-title {{
            font-size: 36px;
            font-weight: bold;
            color: #003366;
            margin-bottom: 20px;
            text-align: center;
        }}

        .slide-content {{
            font-size: 18px;
            line-height: 1.6;
            color: #333;
        }}

        .slide-notes {{
            background: #f9f9f9;
            border-left: 4px solid #003366;
            padding: 15px;
            margin-top: 20px;
            font-style: italic;
            color: #666;
        }}

        .controls {{
            background: white;
            border: 1px solid #ddd;
            border-top: none;
            padding: 20px;
            display: flex;
            gap: 10px;
            flex-wrap: wrap;
        }}

        .btn {{
            padding: 12px 24px;
            border: none;
            border-radius: 6px;
            font-size: 16px;
            cursor: pointer;
            transition: all 0.2s;
        }}

        .btn:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(0,0,0,0.15);
        }}

        .btn-confirm {{
            background: #27ae60;
            color: white;
        }}

        .btn-modify {{
            background: #f39c12;
            color: white;
        }}

        .btn-skip {{
            background: #95a5a6;
            color: white;
        }}

        .btn-redo {{
            background: #e74c3c;
            color: white;
        }}

        .feedback-section {{
            background: white;
            border: 1px solid #ddd;
            border-top: none;
            padding: 20px;
            border-radius: 0 0 8px 8px;
        }}

        .feedback-section h3 {{
            margin-bottom: 15px;
            color: #333;
        }}

        .feedback-input {{
            width: 100%;
            padding: 12px;
            border: 1px solid #ddd;
            border-radius: 6px;
            font-size: 14px;
            resize: vertical;
            min-height: 100px;
        }}

        .feedback-actions {{
            margin-top: 15px;
            display: flex;
            gap: 10px;
        }}

        .status {{
            margin-top: 20px;
            padding: 15px;
            background: #d4edda;
            border: 1px solid #c3e6cb;
            border-radius: 6px;
            color: #155724;
            display: none;
        }}

        .status.show {{
            display: block;
        }}

        .status.error {{
            background: #f8d7da;
            border-color: #f5c6cb;
            color: #721c24;
        }}

        .connection-status {{
            margin-top: 10px;
            padding: 8px 12px;
            border-radius: 4px;
            font-size: 12px;
        }}

        .connection-status.connected {{
            background: #d4edda;
            color: #155724;
        }}

        .connection-status.disconnected {{
            background: #f8d7da;
            color: #721c24;
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>Slide {slide_index} Preview</h1>
            <div class="meta">
                <span>Duration: {duration} seconds</span>
                <div id="connectionStatus" class="connection-status disconnected">
                    未连接
                </div>
            </div>
        </div>

        <div class="slide-preview">
            <div class="slide-title">{title}</div>
            <div class="slide-content">{content}</div>
            {f'<div class="slide-notes"><strong>Notes:</strong> {notes}</div>' if notes else ''}
        </div>

        <div class="controls">
            <button class="btn btn-confirm" onclick="handleAction('confirm')">✅ Confirm</button>
            <button class="btn btn-modify" onclick="handleAction('modify')">✏️ Modify</button>
            <button class="btn btn-skip" onclick="handleAction('skip')">⏭️ Skip</button>
            <button class="btn btn-redo" onclick="handleAction('redo')">🔄 Redo</button>
        </div>

        <div class="feedback-section">
            <h3>Feedback</h3>
            <textarea class="feedback-input" id="feedbackInput"
                placeholder="Enter your modification suggestions..."></textarea>
            <div class="feedback-actions">
                <button class="btn btn-modify" onclick="submitFeedback()">Submit Feedback</button>
                <button class="btn btn-skip" onclick="clearFeedback()">Clear</button>
            </div>
        </div>

        <div class="status" id="status"></div>
    </div>

    <script>
        const SLIDE_INDEX = {slide_index};
        let ws = null;

        // 连接WebSocket服务器
        function connectWebSocket() {{
            const statusDiv = document.getElementById('connectionStatus');

            try {{
                ws = new WebSocket('ws://localhost:8765');

                ws.onopen = function() {{
                    statusDiv.textContent = '已连接';
                    statusDiv.className = 'connection-status connected';
                    console.log('WebSocket connected');
                }};

                ws.onmessage = function(event) {{
                    const data = JSON.parse(event.data);
                    console.log('Received:', data);

                    if (data.type === 'ack') {{
                        showStatus(data.message, false);
                    }} else if (data.type === 'error') {{
                        showStatus(data.message, true);
                    }}
                }};

                ws.onclose = function() {{
                    statusDiv.textContent = '已断开';
                    statusDiv.className = 'connection-status disconnected';
                    console.log('WebSocket disconnected');
                    // 尝试重连
                    setTimeout(connectWebSocket, 3000);
                }};

                ws.onerror = function(error) {{
                    statusDiv.textContent = '连接错误';
                    statusDiv.className = 'connection-status disconnected';
                    console.error('WebSocket error:', error);
                }};
            }} catch (e) {{
                statusDiv.textContent = '连接失败';
                statusDiv.className = 'connection-status disconnected';
            }}
        }}

        // 发送消息到服务器
        function sendMessage(data) {{
            if (ws && ws.readyState === WebSocket.OPEN) {{
                ws.send(JSON.stringify(data));
                return true;
            }} else {{
                showStatus('未连接到服务器，请稍后重试', true);
                return false;
            }}
        }}

        // 处理按钮操作
        function handleAction(action) {{
            const data = {{
                type: 'action',
                action: action,
                slide: SLIDE_INDEX,
                timestamp: new Date().toISOString()
            }};

            if (sendMessage(data)) {{
                showStatus(`Action: ${{action}}`, false);
            }}
        }}

        // 提交反馈
        function submitFeedback() {{
            const feedback = document.getElementById('feedbackInput').value;
            if (!feedback.trim()) {{
                showStatus('Please enter feedback', true);
                return;
            }}

            const data = {{
                type: 'feedback',
                feedback: feedback,
                slide: SLIDE_INDEX,
                timestamp: new Date().toISOString()
            }};

            if (sendMessage(data)) {{
                showStatus('Feedback submitted!', false);
                document.getElementById('feedbackInput').value = '';
            }}
        }}

        // 清空反馈
        function clearFeedback() {{
            document.getElementById('feedbackInput').value = '';
        }}

        // 显示状态信息
        function showStatus(message, isError) {{
            const statusDiv = document.getElementById('status');
            statusDiv.textContent = message;
            statusDiv.className = 'status show' + (isError ? ' error' : '');

            // 3秒后自动隐藏
            setTimeout(() => {{
                statusDiv.className = 'status';
            }}, 3000);
        }}

        // 页面加载时连接WebSocket
        window.onload = connectWebSocket;
    </script>
</body>
</html>"""
        return html

    def get_preview_list(self) -> list:
        """获取所有预览文件列表."""
        previews = []
        for file in self.output_dir.glob("slide_*.html"):
            index = int(file.stem.split("_")[1])
            previews.append({
                "index": index,
                "path": str(file),
                "type": "html"
            })
        return sorted(previews, key=lambda x: x["index"])
