"""修改检测器 - 检测用户在PowerPoint中的手动修改."""

from pptx import Presentation
from pathlib import Path
import json


class ModificationDetector:
    """检测PPTX文件之间的变化."""

    def __init__(self):
        self.changes = []

    def detect_changes(self, original_path: str, modified_path: str) -> list:
        """检测两个PPTX文件之间的变化.

        Args:
            original_path: 原始PPTX文件路径
            modified_path: 修改后PPTX文件路径

        Returns:
            变化列表
        """
        original = Presentation(original_path)
        modified = Presentation(modified_path)

        self.changes = []

        # 比较每一页
        for slide_idx, (orig_slide, mod_slide) in enumerate(
            zip(original.slides, modified.slides)
        ):
            self._compare_slides(orig_slide, mod_slide, slide_idx)

        return self.changes

    def _compare_slides(self, orig_slide, mod_slide, slide_index: int):
        """比较两张幻灯片."""
        # 比较每个形状
        orig_shapes = list(orig_slide.shapes)
        mod_shapes = list(mod_slide.shapes)

        # 检查形状数量变化
        if len(orig_shapes) != len(mod_shapes):
            self.changes.append({
                "type": "shape_count",
                "slide": slide_index,
                "original": len(orig_shapes),
                "modified": len(mod_shapes)
            })

        # 比较共同形状
        for orig_shape, mod_shape in zip(orig_shapes, mod_shapes):
            self._compare_shapes(orig_shape, mod_shape, slide_index)

    def _compare_shapes(self, orig_shape, mod_shape, slide_index: int):
        """比较两个形状."""
        shape_name = orig_shape.name

        # 检测文本变化
        if orig_shape.has_text_frame and mod_shape.has_text_frame:
            orig_text = orig_shape.text_frame.text
            mod_text = mod_shape.text_frame.text

            if orig_text != mod_text:
                self.changes.append({
                    "type": "text",
                    "slide": slide_index,
                    "shape": shape_name,
                    "original": orig_text,
                    "modified": mod_text
                })

        # 检测位置变化
        if (orig_shape.left != mod_shape.left or
            orig_shape.top != mod_shape.top):
            self.changes.append({
                "type": "position",
                "slide": slide_index,
                "shape": shape_name,
                "original": {
                    "left": orig_shape.left,
                    "top": orig_shape.top
                },
                "modified": {
                    "left": mod_shape.left,
                    "top": mod_shape.top
                }
            })

        # 检测大小变化
        if (orig_shape.width != mod_shape.width or
            orig_shape.height != mod_shape.height):
            self.changes.append({
                "type": "size",
                "slide": slide_index,
                "shape": shape_name,
                "original": {
                    "width": orig_shape.width,
                    "height": orig_shape.height
                },
                "modified": {
                    "width": mod_shape.width,
                    "height": mod_shape.height
                }
            })

        # 检测颜色变化（如果形状有填充）
        try:
            if hasattr(orig_shape, 'fill') and hasattr(mod_shape, 'fill'):
                orig_fill = orig_shape.fill
                mod_fill = mod_shape.fill

                if orig_fill.type is not None and mod_fill.type is not None:
                    if hasattr(orig_fill, 'fore_color') and hasattr(mod_fill, 'fore_color'):
                        orig_color = orig_fill.fore_color.rgb
                        mod_color = mod_fill.fore_color.rgb

                        if orig_color != mod_color:
                            self.changes.append({
                                "type": "color",
                                "slide": slide_index,
                                "shape": shape_name,
                                "original": str(orig_color),
                                "modified": str(mod_color)
                            })
        except Exception:
            # 颜色检测可能失败，忽略
            pass

    def get_changes_summary(self) -> dict:
        """获取变化摘要."""
        summary = {
            "total_changes": len(self.changes),
            "by_type": {},
            "by_slide": {}
        }

        for change in self.changes:
            # 按类型统计
            change_type = change["type"]
            if change_type not in summary["by_type"]:
                summary["by_type"][change_type] = 0
            summary["by_type"][change_type] += 1

            # 按页面统计
            slide_idx = change["slide"]
            if slide_idx not in summary["by_slide"]:
                summary["by_slide"][slide_idx] = 0
            summary["by_slide"][slide_idx] += 1

        return summary

    def export_changes(self, output_path: str):
        """导出变化到JSON文件."""
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(self.changes, f, indent=2, ensure_ascii=False)

    def has_significant_changes(self) -> bool:
        """检查是否有显著变化."""
        # 检查是否有文本、位置、大小或颜色变化
        significant_types = {"text", "position", "size", "color"}
        return any(change["type"] in significant_types for change in self.changes)
