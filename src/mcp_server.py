"""SCI-two-ppt MCP Server entry point.

Exposes 11 tools for Claude Code to orchestrate the paper-to-PPT workflow.
"""

from mcp.server.fastmcp import FastMCP
import json
import asyncio
from pathlib import Path

from src.core.config import load_config
from src.core.workspace import Workspace
from src.parsers.paper_parser import parse_papers as _parse_papers
from src.parsers.figure_extractor import extract_figures as _extract_figures
from src.generators.slide_builder import build_slide as _build_slide
from src.styles.academic_styles import get_academic_style as _get_academic_style
from src.styles.slide_templates import get_slide_template as _get_slide_template
from src.styles.citation_formats import get_citation_format as _get_citation_format

mcp = FastMCP("sci-two-ppt")
config = load_config()


@mcp.tool()
def parse_papers(papers: list[str], workspace_path: str = "") -> str:
    """解析论文，提取原始文本（PDF/Word）。

    Input: papers - 论文文件路径列表（支持 .pdf 和 .docx）
    Output: JSON {
        "papers": [{
            "path": str,
            "format": "pdf|docx",
            "raw_text": str,
            "page_count": int (PDF) | "paragraph_count": int (Word),
            "word_count": int,
            "char_count": int
        }],
        "errors": [str],
        "total": int,
        "failed": int
    }

    注意：此工具只提取原始文本，不做结构分析。
    结构分析（标题、摘要、关键发现等）应由子Agent完成。
    """
    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        result = asyncio.run(_parse_papers(papers))
        return json.dumps(result, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def extract_figures(pdf_path: str, output_dir: str = "", workspace_path: str = "") -> str:
    """提取论文PDF中的图表和图片。

    Input: pdf_path - PDF文件路径, output_dir - 图片输出目录
    Output: JSON {
        "figures": [{"path": str, "caption": str, "type": "figure|table", "page": int}]
    }
    """
    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    output_dir = output_dir or str(ws.path / "assets" / "figures")

    try:
        figures = asyncio.run(_extract_figures(pdf_path, output_dir))
        return json.dumps({"figures": figures}, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def build_goal(paper_analysis: str, requirements: str, workspace_path: str = "") -> str:
    """构建结构化的PPT目标文档。

    Input: paper_analysis - parse_papers输出的JSON或analysis.json路径, requirements - 用户需求文本或requirements.md路径
    Output: JSON {
        "goal_content": str,
        "goal_path": str,
        "sections": [str],
        "slide_count_estimate": int
    }

    读取论文分析和用户需求，按标准模板生成 goal.md 并保存到 workspace。
    """
    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        # 加载分析数据
        if Path(paper_analysis).exists():
            with open(paper_analysis, 'r', encoding='utf-8') as f:
                analysis = json.load(f)
            analysis_text = json.dumps(analysis, ensure_ascii=False, indent=2)
        else:
            analysis_text = paper_analysis
            analysis = {}

        # 加载需求
        if Path(requirements).exists():
            with open(requirements, 'r', encoding='utf-8') as f:
                requirements_text = f.read()
        else:
            requirements_text = requirements

        title = analysis.get("title", "未指定标题")
        innovations = analysis.get("innovations", [])
        key_findings = analysis.get("key_findings", [])
        research_field = analysis.get("research_field", "未指定领域")
        conclusions = analysis.get("conclusions", "")

        # 从需求中提取配置
        duration = "10分钟"
        audience = "学术听众"
        purpose = "学术会议汇报"
        for line in requirements_text.split('\n'):
            if '时长' in line or 'duration' in line.lower():
                duration = line.split('：')[-1].split(':')[-1].strip() or duration
            if '听众' in line or 'audience' in line.lower():
                audience = line.split('：')[-1].split(':')[-1].strip() or audience
            if '用途' in line or 'purpose' in line.lower():
                purpose = line.split('：')[-1].split(':')[-1].strip() or purpose

        innovations_text = '\n'.join(f"- {inn}" for inn in innovations) if innovations else "- 待分析"
        findings_text = '\n'.join(f"- {f}" for f in key_findings) if key_findings else "- 待分析"

        goal_content = f"""# PPT目标文档

## 1. PPT概述
- 标题：{title}
- 用途：{purpose}
- 时长：{duration}
- 听众：{audience}
- 研究领域：{research_field}

## 2. 内容结构
1. 封面（30秒）
2. 研究背景（1分钟）
3. 研究目的与创新点（1分钟）
4. 研究方法（1.5分钟）
5. 实验结果（2分钟）
6. 结论与展望（1分钟）
7. 致谢（30秒）

## 3. 核心要点
### 创新点
{innovations_text}

### 关键发现
{findings_text}

## 4. 视觉方向（高层建议，详细规范由 design_spec.md 定义）
- 配色倾向：根据研究领域自动选择
- 风格倾向：学术专业风格

## 5. 讲解策略
- 开场白：介绍研究背景和动机
- 重点强调：创新点和关键数据
- 过渡语：自然衔接各章节

## 6. 信息补全
- 结论：{conclusions}
- 待补充：请根据审查结果补充

---

> 此文档由 build_goal 工具自动生成，请根据实际情况修改。
"""
        goal_path = ws.save_artifact("goal.md", goal_content)

        return json.dumps({
            "goal_content": goal_content,
            "goal_path": str(goal_path),
            "sections": ["概述", "内容结构", "核心要点", "视觉方向", "讲解策略", "信息补全"],
            "slide_count_estimate": 7,
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def run_subagent(
    agent_type: str,
    goal: str,
    context: str = "{}",
    workspace_path: str = "",
) -> str:
    """执行子Agent任务。

    Input: agent_type - Agent类型（见下）, goal - goal.md内容, context - 上下文JSON
           支持的 agent_type:
             - paper_keypoints: 论文要点提取 (Step5 Agent1)
             - innovation_points: 核心创新点提炼 (Step5 Agent2)
             - simulation_code: 仿真代码分析 (Step5 Agent3)
             - visual_resources: 学术配图搜集 (Step5 Agent4)
             - ui_design: UI风格设计 (Step5 Agent5)
             - chapter_structure: 章节结构安排 (Step5 Agent6)
             - speaker_notes: 讲解备注 (Step5 Agent7)
    Output: JSON {
        "agent_type": str,
        "result_md": str,
        "output_path": str,
        "assets": [str]
    }
    """
    valid_types = [
        "paper_keypoints", "innovation_points", "simulation_code",
        "visual_resources", "ui_design", "chapter_structure", "speaker_notes",
    ]
    if agent_type not in valid_types:
        return json.dumps({"error": f"Invalid agent_type: {agent_type}. Must be one of {valid_types}"})

    ws = Workspace(workspace_path or config.default_workspace)
    ws_path = ws.ensure_exists()

    # 读取分析数据
    analysis_path = ws_path / "papers" / "analysis.json"
    analysis = {}
    if analysis_path.exists():
        with open(analysis_path, 'r', encoding='utf-8') as f:
            analysis = json.load(f)

    # 读取已有产出（用于依赖分析）
    agent_results_dir = ws_path / "agent_results"
    agent_results_dir.mkdir(parents=True, exist_ok=True)

    def _read_agent_result(filename):
        p = agent_results_dir / filename
        if p.exists():
            return p.read_text(encoding='utf-8')
        return ""

    output_filename = ""
    result_md = ""

    if agent_type == "paper_keypoints":
        innovations = analysis.get("innovations", [])
        key_findings = analysis.get("key_findings", [])
        methods = analysis.get("methods", "")
        results_text = analysis.get("results", "")
        result_md = f"""# 论文要点提取

## 核心发现
{chr(10).join(f'- {f}' for f in key_findings) if key_findings else '- 待分析'}

## 创新点
{chr(10).join(f'- {i}' for i in innovations) if innovations else '- 待分析'}

## 研究方法
{methods[:800] if methods else '待分析'}

## 实验结果
{results_text[:800] if results_text else '待分析'}
"""
        output_filename = "01_paper_keypoints.md"

    elif agent_type == "innovation_points":
        innovations = analysis.get("innovations", [])
        result_md = f"""# 核心创新点提炼

## 创新点分析
{chr(10).join(f'### 创新点 {i+1}\n{inn}' for i, inn in enumerate(innovations)) if innovations else '### 待分析\n基于论文内容提炼创新点'}

## 与现有研究对比
- 对比分析待补充

## 技术优势
- 技术优势待补充
"""
        output_filename = "02_innovation_points.md"

    elif agent_type == "simulation_code":
        methods = analysis.get("methods", "")
        results_text = analysis.get("results", "")
        result_md = f"""# 仿真代码分析

## 研究方法概述
{methods[:1000] if methods else '待分析'}

## 关键参数
- 参数分析待补充

## 仿真结果
{results_text[:1000] if results_text else '待分析'}
"""
        output_filename = "03_simulation_code.md"

    elif agent_type == "visual_resources":
        research_field = analysis.get("research_field", "未知领域")
        core_keywords = analysis.get("core_keywords", [])
        result_md = f"""# 学术配图搜集

## 研究领域
{research_field}

## 核心关键词
{chr(10).join(f'- {k}' for k in core_keywords) if core_keywords else '- 待补充'}

## 搜集计划
1. 从论文中提取已有图表
2. 搜索相关领域的示意图
3. 生成数据可视化图表

## 素材清单
- 素材待搜集
"""
        output_filename = "04_visual_resources.md"

    elif agent_type == "ui_design":
        goal_path = ws_path / "goal.md"
        goal_content = goal_path.read_text(encoding='utf-8') if goal_path.exists() else ""
        design_spec_path = ws_path / "design_spec.md"
        design_spec = design_spec_path.read_text(encoding='utf-8') if design_spec_path.exists() else ""

        domain = analysis.get("research_field", "general").lower()
        primary_color, secondary_color, accent_color = "#1A5276", "#2E86C1", "#E74C3C"
        if "optics" in domain or "光学" in domain:
            primary_color, secondary_color, accent_color = "#003366", "#6699CC", "#FF6600"
        elif "bio" in domain or "生物" in domain:
            primary_color, secondary_color, accent_color = "#196F3D", "#58D68D", "#E74C3C"
        elif "computer" in domain or "计算机" in domain:
            primary_color, secondary_color, accent_color = "#2C3E50", "#9B59B6", "#F39C12"

        result_md = f"""# UI风格设计

## 配色方案
- 主色：{primary_color}（标题、边框）
- 辅色：{secondary_color}（副标题、背景）
- 强调色：{accent_color}（重点标注）
- 背景色：#FFFFFF

## 字体规范
- 标题：Arial Bold 36pt
- 副标题：Arial Semibold 24pt
- 正文：Arial 18pt
- 公式：Cambria Math 20pt
- 图注：Arial 14pt

## 页面布局
- 封面页：居中布局，渐变背景（主色→辅色）
- 内容页：上标题栏 + 左文右图
- 图表页：标题 + 居中图表 + 图注
- 总结页：居中布局，渐变背景

## 设计决策来源
{design_spec[:500] if design_spec else '（待生成 design_spec.md）'}
"""
        output_filename = "05_ui_design.md"

        # 同时生成 spec_lock.md
        spec_lock = f"""# 视觉设计执行锁
colors:
  primary: "{primary_color}"
  secondary: "{secondary_color}"
  accent: "{accent_color}"
  background: "#FFFFFF"
  text_primary: "#333333"
  text_secondary: "#666666"
  text_light: "#FFFFFF"
fonts:
  title:
    family: "Arial"
    size: 36
    weight: "bold"
    color: "{primary_color}"
  subtitle:
    family: "Arial"
    size: 24
    weight: "semibold"
    color: "{secondary_color}"
  body:
    family: "Arial"
    size: 18
    weight: "normal"
    color: "#333333"
  caption:
    family: "Arial"
    size: 14
    weight: "normal"
    color: "#666666"
layouts:
  title_slide:
    background_type: "gradient"
    background_colors: ["{primary_color}", "{secondary_color}"]
  content_slide:
    title_position: "top"
    content_position: "left"
    figure_position: "right"
  chart_slide:
    title_position: "top"
    chart_position: "center"
  conclusion_slide:
    background_type: "gradient"
    background_colors: ["{primary_color}", "{secondary_color}"]
"""
        spec_lock_path = ws_path / "papers" / "spec_lock.md"
        spec_lock_path.parent.mkdir(parents=True, exist_ok=True)
        spec_lock_path.write_text(spec_lock, encoding='utf-8')

    elif agent_type == "chapter_structure":
        result_md = """# 章节结构安排

## PPT结构
1. **封面**（30秒）- 标题、作者、单位
2. **研究背景**（1分钟）- 问题引出、研究意义
3. **研究目的与创新点**（1分钟）- 核心创新点展示
4. **研究方法**（1.5分钟）- 方法原理、实验设计
5. **实验结果**（2分钟）- 关键数据、图表展示
6. **结论与展望**（1分钟）- 核心结论、未来工作
7. **致谢**（30秒）- 感谢语

## 时间分配依据
基于创新点和关键发现的重要性分配时间。

## 内容逻辑
- 从背景到问题到方案到结果到结论
- 创新点在第3页重点展示
"""
        output_filename = "06_chapter_structure.md"

    elif agent_type == "speaker_notes":
        result_md = """# 讲解备注

## 开场白
各位老师、同学好，今天汇报的题目是关于本研究的工作。

## 各页备注

### 封面（Page 1）
- 介绍论文标题和作者信息

### 研究背景（Page 2）
- 介绍研究问题的背景和意义

### 研究目的与创新点（Page 3）
- 明确研究目标，重点强调核心创新点

### 研究方法（Page 4）
- 介绍方法原理，说明实验设计

### 实验结果（Page 5）
- 展示关键数据，分析图表结果

### 结论与展望（Page 6）
- 总结核心结论，展望未来工作

### 致谢（Page 7）
- 感谢导师和合作者

## 过渡语设计
- 章节之间："接下来介绍..."
- 重点之前："特别需要注意的是..."
- 数据展示："从数据可以看出..."
"""
        output_filename = "07_speaker_notes.md"

    # 保存到文件
    output_path = agent_results_dir / output_filename
    output_path.write_text(result_md, encoding='utf-8')

    return json.dumps({
        "agent_type": agent_type,
        "result_md": result_md,
        "output_path": str(output_path),
        "assets": [],
    }, ensure_ascii=False, indent=2)


@mcp.tool()
def read_pptx(pptx_path: str, workspace_path: str = "") -> str:
    """读取PPTX文件状态，获取幻灯片信息。

    Input: pptx_path - PPTX文件路径
    Output: JSON {
        "slide_count": int,
        "slides": [{"index": int, "shapes": [{"name": str, "type": str, "text": str}]}]
    }
    """
    from pptx import Presentation

    try:
        prs = Presentation(pptx_path)
        slides_info = []

        for idx, slide in enumerate(prs.slides):
            shapes_info = []
            for shape in slide.shapes:
                shape_info = {
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "text": shape.text if shape.has_text_frame else ""
                }
                shapes_info.append(shape_info)

            slides_info.append({
                "index": idx,
                "shapes": shapes_info
            })

        return json.dumps({
            "slide_count": len(prs.slides),
            "slides": slides_info
        }, ensure_ascii=False, indent=2)

    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def diff_pptx(original: str, modified: str, workspace_path: str = "") -> str:
    """对比两版PPTX的差异，检测用户手动修改。

    Input: original - 原版PPTX路径, modified - 修改后PPTX路径
    Output: JSON {
        "changes": [{"slide_index": int, "type": str, "detail": str}],
        "summary": str
    }
    """
    from pptx import Presentation

    try:
        prs_orig = Presentation(original)
        prs_mod = Presentation(modified)

        changes = []

        # 比较幻灯片数量
        if len(prs_orig.slides) != len(prs_mod.slides):
            changes.append({
                "slide_index": -1,
                "type": "slide_count",
                "detail": f"幻灯片数量从 {len(prs_orig.slides)} 变为 {len(prs_mod.slides)}"
            })

        # 逐页比较
        for idx in range(min(len(prs_orig.slides), len(prs_mod.slides))):
            slide_orig = prs_orig.slides[idx]
            slide_mod = prs_mod.slides[idx]

            # 比较文本内容
            for shape_orig, shape_mod in zip(slide_orig.shapes, slide_mod.shapes):
                if shape_orig.has_text_frame and shape_mod.has_text_frame:
                    if shape_orig.text != shape_mod.text:
                        changes.append({
                            "slide_index": idx,
                            "type": "text",
                            "detail": f"形状 '{shape_orig.name}' 文本已修改"
                        })

                # 比较位置和大小
                if (shape_orig.left != shape_mod.left or
                    shape_orig.top != shape_mod.top or
                    shape_orig.width != shape_mod.width or
                    shape_orig.height != shape_mod.height):
                    changes.append({
                        "slide_index": idx,
                        "type": "position",
                        "detail": f"形状 '{shape_orig.name}' 位置/大小已修改"
                    })

        summary = f"共检测到 {len(changes)} 处变化"
        if not changes:
            summary = "未检测到明显变化"

        return json.dumps({
            "changes": changes,
            "summary": summary
        }, ensure_ascii=False, indent=2)

    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def generate_blueprint(goal: str, agent_results: str, workspace_path: str = "") -> str:
    """生成详细的PPT蓝图。

    Input: goal - goal.md内容或路径, agent_results - Agent结果JSON或目录路径
    Output: JSON {
        "blueprint_yaml": str,
        "blueprint_path": str,
        "slide_count": int
    }

    读取 goal.md 和 agent_results，根据章节结构和内容生成 YAML 蓝图。
    """
    ws = Workspace(workspace_path or config.default_workspace)
    ws_path = ws.ensure_exists()

    try:
        # 读取 goal.md
        if Path(goal).exists():
            with open(goal, 'r', encoding='utf-8') as f:
                goal_content = f.read()
        else:
            goal_content = goal

        # 读取 agent_results 目录
        agent_dir = Path(agent_results) if Path(agent_results).is_dir() else ws_path / "agent_results"
        agent_texts = {}
        if agent_dir.exists():
            for p in sorted(agent_dir.glob("*.md")):
                agent_texts[p.stem] = p.read_text(encoding='utf-8')

        # 从 goal.md 提取标题
        title = "论文标题"
        for line in goal_content.split('\n'):
            if line.strip().startswith('- 标题：') or line.strip().startswith('- 标题:'):
                title = line.split('：')[-1].split(':')[-1].strip()
                break

        # 从 chapter_structure 提取章节信息
        chapter_text = agent_texts.get("06_chapter_structure", "")
        speaker_text = agent_texts.get("07_speaker_notes", "")

        # 构建蓝图
        slides = []
        slide_index = 0

        # Slide 0: 封面
        slides.append({
            "index": slide_index,
            "type": "title",
            "title": title,
            "subtitle": "作者信息 | 单位信息",
            "notes": "介绍论文标题和作者信息",
            "duration_seconds": 30,
            "layout": "centered",
        })
        slide_index += 1

        # 标准章节配置
        standard_sections = [
            ("研究背景", "content", "介绍研究问题的背景和意义", 60, "left_text_right_image"),
            ("研究目的与创新点", "content", "明确研究目标，重点强调核心创新点", 60, "left_text_right_image"),
            ("研究方法", "content", "介绍方法原理，说明实验设计", 90, "left_text_right_image"),
            ("实验结果", "chart", "展示关键数据，分析图表结果", 120, "chart_center"),
            ("结论与展望", "conclusion", "总结核心结论，展望未来工作", 60, "centered"),
            ("致谢", "conclusion", "感谢导师和合作者", 30, "centered"),
        ]

        # 从 chapter_structure 解析实际章节（如果有）
        actual_sections = []
        for line in chapter_text.split('\n'):
            line = line.strip()
            if line.startswith(('1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.')):
                # 提取章节名和时间
                parts = line.split('**')
                if len(parts) >= 3:
                    section_name = parts[1].strip()
                    # 提取时间
                    time_part = line.split('（')[-1].split('）')[0] if '（' in line else "60秒"
                    try:
                        if '分钟' in time_part:
                            duration = int(float(time_part.replace('分钟', '').strip()) * 60)
                        elif '秒' in time_part:
                            duration = int(time_part.replace('秒', '').strip())
                        else:
                            duration = 60
                    except ValueError:
                        duration = 60

                    slide_type = "title" if "封面" in section_name else \
                                 "chart" if "结果" in section_name or "实验" in section_name else \
                                 "conclusion" if "结论" in section_name or "致谢" in section_name else "content"

                    actual_sections.append({
                        "name": section_name,
                        "type": slide_type,
                        "duration": duration,
                    })

        # 如果解析到实际章节则用实际的，否则用标准的
        sections_to_use = actual_sections if len(actual_sections) >= 3 else [
            {"name": s[0], "type": s[1], "duration": s[4]} for s in standard_sections
        ]

        for section in sections_to_use:
            slide_def = {
                "index": slide_index,
                "type": section["type"],
                "title": section["name"],
                "content": f"在此展示{section['name']}的相关内容",
                "notes": f"讲解{section['name']}",
                "duration_seconds": section.get("duration", 60),
                "layout": "chart_center" if section["type"] == "chart" else "left_text_right_image",
            }
            slides.append(slide_def)
            slide_index += 1

        # 构建 YAML
        yaml_lines = ["slides:"]
        for s in slides:
            yaml_lines.append(f"  - index: {s['index']}")
            yaml_lines.append(f"    type: {s['type']}")
            yaml_lines.append(f"    title: \"{s['title']}\"")
            if "subtitle" in s:
                yaml_lines.append(f"    subtitle: \"{s['subtitle']}\"")
            if "content" in s:
                yaml_lines.append(f"    content: \"{s['content']}\"")
            yaml_lines.append(f"    notes: \"{s['notes']}\"")
            yaml_lines.append(f"    duration_seconds: {s['duration_seconds']}")
            yaml_lines.append(f"    layout: \"{s.get('layout', 'left_text_right_image')}\"")
            yaml_lines.append("")

        blueprint_yaml = "\n".join(yaml_lines)

        # 保存到文件
        blueprint_path = ws_path / "blueprint.yaml"
        blueprint_path.write_text(blueprint_yaml, encoding='utf-8')

        return json.dumps({
            "blueprint_yaml": blueprint_yaml,
            "blueprint_path": str(blueprint_path),
            "slide_count": len(slides),
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def build_slide(
    blueprint: str,
    slide_index: int,
    modifications: str = "",
    workspace_path: str = "",
) -> str:
    """根据蓝图生成单页幻灯片。

    Input: blueprint - 蓝图YAML, slide_index - 页码(从0开始), modifications - 修改意见
    Output: JSON {
        "slide_index": int,
        "pptx_path": str
    }
    """
    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        result = asyncio.run(_build_slide(blueprint, slide_index, modifications))
        return json.dumps(result, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def render_preview(pptx_path: str, slide_index: int, workspace_path: str = "") -> str:
    """渲染幻灯片为HTML预览。

    Input: pptx_path - PPTX文件路径, slide_index - 页码
    Output: JSON {
        "html_path": str,
        "pptx_path": str
    }

    从PPTX文件生成HTML预览，保持与PowerPoint一致的视觉效果。
    """
    from src.generators.preview.preview_generator import PreviewGenerator

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        if not Path(pptx_path).exists():
            return json.dumps({"error": f"PPTX文件不存在: {pptx_path}"}, ensure_ascii=False, indent=2)

        preview_gen = PreviewGenerator(str(ws.path / "preview"))
        html_path = preview_gen.generate_html_from_pptx(pptx_path, slide_index)

        return json.dumps({
            "html_path": html_path,
            "pptx_path": pptx_path,
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def generate_pptx(
    blueprint: str,
    slide_dir: str = "",
    output_path: str = "",
    workspace_path: str = "",
) -> str:
    """最终打包生成PPTX文件。

    Input: blueprint - 蓝图YAML或路径, slide_dir - 单页PPTX目录, output_path - 输出路径
    Output: JSON {
        "pptx_path": str,
        "report_md": str,
        "slide_count": int
    }

    有两种模式：
    1. 如果 slide_dir 中有单页 PPTX 文件，合并它们
    2. 否则根据 blueprint 生成新 PPTX（使用 SVG 引擎或 python-pptx）
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    import yaml as yaml_lib

    ws = Workspace(workspace_path or config.default_workspace)
    ws_path = ws.ensure_exists()
    output = output_path or str(ws_path / "output" / "output_final.pptx")
    Path(output).parent.mkdir(parents=True, exist_ok=True)

    try:
        # 读取蓝图
        if Path(blueprint).exists():
            with open(blueprint, 'r', encoding='utf-8') as f:
                blueprint_text = f.read()
        else:
            blueprint_text = blueprint

        # 尝试解析蓝图
        try:
            bp = yaml_lib.safe_load(blueprint_text)
            slides_def = bp.get("slides", []) if bp else []
        except Exception:
            slides_def = []

        # 检查是否有单页 PPTX 文件可以合并
        slide_dir_path = Path(slide_dir) if slide_dir else ws_path / "preview"
        pptx_files = sorted(slide_dir_path.glob("slide_*.pptx")) if slide_dir_path.exists() else []

        if pptx_files:
            # 模式1：合并已有单页 PPTX
            merged_prs = Presentation()
            # 设置为16:9
            merged_prs.slide_width = Inches(13.333)
            merged_prs.slide_height = Inches(7.5)

            for pptx_file in pptx_files:
                try:
                    src_prs = Presentation(str(pptx_file))
                    for slide in src_prs.slides:
                        # 复制幻灯片
                        layout = merged_prs.slide_layouts[6]  # 空白布局
                        new_slide = merged_prs.slides.add_slide(layout)

                        # 复制所有形状
                        for shape in slide.shapes:
                            el = shape._element
                            new_slide.shapes._spTree.append(el)
                except Exception as e:
                    continue

            merged_prs.save(output)
            slide_count = len(merged_prs.slides)

        elif slides_def:
            # 模式2：根据蓝图创建新 PPTX
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            for slide_def in slides_def:
                slide_type = slide_def.get("type", "content")
                title = slide_def.get("title", "")
                content = slide_def.get("content", "")
                notes = slide_def.get("notes", "")
                subtitle = slide_def.get("subtitle", "")

                layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(layout)

                # 添加标题
                if title:
                    from pptx.util import Inches as In
                    txBox = slide.shapes.add_textbox(In(0.8), In(0.4), In(11.7), In(1.0))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = title
                    p.font.size = Pt(36)
                    p.font.bold = True

                # 添加副标题（封面页）
                if subtitle and slide_type == "title":
                    txBox2 = slide.shapes.add_textbox(In(2.0), In(3.0), In(9.3), In(1.0))
                    tf2 = txBox2.text_frame
                    p2 = tf2.paragraphs[0]
                    p2.text = subtitle
                    p2.font.size = Pt(24)

                # 添加内容
                if content and slide_type != "title":
                    txBox3 = slide.shapes.add_textbox(In(0.8), In(1.8), In(11.7), In(5.0))
                    tf3 = txBox3.text_frame
                    tf3.word_wrap = True
                    p3 = tf3.paragraphs[0]
                    p3.text = content
                    p3.font.size = Pt(18)

                # 添加备注
                if notes:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes

            prs.save(output)
            slide_count = len(prs.slides)
        else:
            # 无蓝图无单页，创建空PPTX
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            prs.save(output)
            slide_count = 0

        # 生成制作报告
        report = f"""# PPT制作报告

## 基本信息
- 输出文件：{output}
- 总页数：{slide_count}
- 蓝图来源：{'蓝图文件' if Path(blueprint).exists() else '直接输入'}

## 生成方式
{'合并单页PPTX文件' if pptx_files else '根据蓝图生成'}

## 文件清单
- {output} - 最终PPT
"""
        report_path = ws_path / "production_report.md"
        report_path.write_text(report, encoding='utf-8')

        return json.dumps({
            "pptx_path": output,
            "report_md": report,
            "slide_count": slide_count,
            "report_path": str(report_path),
        }, ensure_ascii=False, indent=2)

    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def get_academic_style(domain: str = "general", workspace_path: str = "") -> str:
    """获取学术PPT规范（配色、字体、排版）。

    Input: domain - 领域(optics/physics/chemistry/computer_science/biology/general)
    Output: JSON with color scheme, fonts, and margin specifications
    """
    try:
        result = asyncio.run(_get_academic_style(domain))
        return json.dumps(result, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def get_slide_template(slide_type: str, workspace_path: str = "") -> str:
    """获取页面模板定义。

    Input: slide_type - 页面类型(title/content/two_column/chart/image/comparison/bullet_points/conclusion/thank_you)
    Output: JSON with template layout and element specifications
    """
    try:
        result = asyncio.run(_get_slide_template(slide_type))
        return json.dumps(result, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def get_citation_format(format_type: str = "IEEE", workspace_path: str = "") -> str:
    """获取引用格式规范。

    Input: format_type - 格式(IEEE/APA/MLA/Chicago/Harvard)
    Output: JSON with citation format specifications and examples
    """
    try:
        result = asyncio.run(_get_citation_format(format_type))
        return json.dumps(result, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def cleanup_workspace(workspace_path: str = "") -> str:
    """整理清理工作空间，移除中间文件。

    Input: workspace_path - 工作空间路径
    Output: JSON {
        "cleaned_files": [str],
        "kept_files": [str],
        "report": str
    }
    """
    ws = Workspace(workspace_path or config.default_workspace)

    if not ws.path.exists():
        return json.dumps({"report": "工作空间不存在"}, ensure_ascii=False, indent=2)

    cleaned = []
    kept = []

    # 清理中间文件
    intermediate_dirs = ["agent_results", "preview"]
    for dir_name in intermediate_dirs:
        dir_path = ws.path / dir_name
        if dir_path.exists():
            for f in dir_path.rglob("*"):
                if f.is_file():
                    cleaned.append(str(f.relative_to(ws.path)))
                    f.unlink()

    # 保留重要文件
    important_files = [
        "output/output_final.pptx",
        "blueprint.yaml",
        "goal.md",
        "requirements.md",
        "design_spec.md",
        "production_report.md",
    ]
    for file_name in important_files:
        file_path = ws.path / file_name
        if file_path.exists():
            kept.append(file_name)

    report = f"清理完成: 删除 {len(cleaned)} 个文件, 保留 {len(kept)} 个重要文件"

    return json.dumps({
        "cleaned_files": cleaned[:50],  # 限制显示数量
        "kept_files": kept,
        "report": report
    }, ensure_ascii=False, indent=2)


@mcp.tool()
def generate_preview(slide_data: str, slide_index: int, workspace_path: str = "") -> str:
    """生成幻灯片预览（HTML/图片/PPT）。

    Input: slide_data - 幻灯片数据JSON, slide_index - 幻灯片索引
    Output: JSON {
        "html_path": str,
        "image_path": str,
        "pptx_path": str
    }
    """
    from src.generators.preview.preview_generator import PreviewGenerator

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        preview_gen = PreviewGenerator(str(ws.path / "preview"))
        slide_dict = json.loads(slide_data)

        html_path = preview_gen.generate_html_preview(slide_dict, slide_index)

        return json.dumps({
            "html_path": html_path,
            "image_path": "",
            "pptx_path": ""
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def generate_preview_from_pptx(pptx_path: str, slide_index: int, workspace_path: str = "") -> str:
    """从PPTX文件生成HTML预览（保持与PowerPoint一致）。

    Input: pptx_path - PPTX文件路径, slide_index - 幻灯片索引
    Output: JSON {
        "html_path": str,
        "pptx_path": str
    }
    """
    from src.generators.preview.preview_generator import PreviewGenerator

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        preview_gen = PreviewGenerator(str(ws.path / "preview"))
        html_path = preview_gen.generate_html_from_pptx(pptx_path, slide_index)

        return json.dumps({
            "html_path": html_path,
            "pptx_path": pptx_path
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def generate_svg_slide(
    title: str,
    content: str,
    layout: str = "content",
    slide_index: int = 0,
    workspace_path: str = ""
) -> str:
    """生成单页幻灯片的SVG代码。

    Input: title - 标题, content - 内容, layout - 布局类型(title/content/chart/conclusion), slide_index - 页码
    Output: JSON {
        "svg_content": str,
        "svg_path": str
    }
    """
    from src.svg_engine.svg_generator import SvgGenerator

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        generator = SvgGenerator()
        svg_content = generator.generate_slide_svg(title, content, layout)
        svg_path = generator.save_svg(svg_content, f"slide_{slide_index}.svg")

        return json.dumps({
            "svg_content": svg_content,
            "svg_path": svg_path
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def generate_svg_chart(
    chart_type: str,
    title: str,
    values: str,
    labels: str = "",
    slide_index: int = 0,
    workspace_path: str = ""
) -> str:
    """生成图表SVG代码。

    Input: chart_type - 图表类型(bar/pie/line), title - 标题, values - 数据值(JSON数组), labels - 标签(JSON数组)
    Output: JSON {
        "svg_content": str,
        "svg_path": str
    }
    """
    from src.svg_engine.svg_generator import SvgGenerator

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        generator = SvgGenerator()
        values_list = json.loads(values)
        labels_list = json.loads(labels) if labels else None

        if chart_type == "bar":
            svg_content = generator.generate_bar_chart_svg(title, values_list, labels_list)
        elif chart_type == "pie":
            svg_content = generator.generate_pie_chart_svg(title, values_list, labels_list)
        else:
            return json.dumps({"error": f"不支持的图表类型: {chart_type}"}, ensure_ascii=False, indent=2)

        svg_path = generator.save_svg(svg_content, f"chart_{slide_index}.svg")

        return json.dumps({
            "svg_content": svg_content,
            "svg_path": svg_path
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def svg_to_pptx(svg_path: str, output_path: str, workspace_path: str = "") -> str:
    """将SVG文件转换为PPTX文件。

    Input: svg_path - SVG文件路径, output_path - 输出PPTX路径
    Output: JSON {
        "success": bool,
        "pptx_path": str,
        "errors": list
    }
    """
    from src.svg_engine.svg_converter import SvgToPptxConverter

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        converter = SvgToPptxConverter(str(ws.path / "output"))
        result = converter.convert_svg_file_to_pptx(svg_path, output_path)

        return json.dumps(result, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def check_svg_quality(svg_path: str, workspace_path: str = "") -> str:
    """检查SVG文件质量。

    Input: svg_path - SVG文件路径
    Output: JSON {
        "passed": bool,
        "errors": list,
        "warnings": list
    }
    """
    from src.svg_engine.svg_converter import SvgQualityChecker

    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            svg_content = f.read()

        checker = SvgQualityChecker()
        result = checker.check_svg(svg_content)

        return json.dumps(result, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def detect_modifications(original: str, modified: str, workspace_path: str = "") -> str:
    """检测PPTX文件之间的变化。

    Input: original - 原始PPTX路径, modified - 修改后PPTX路径
    Output: JSON {
        "changes": list,
        "summary": dict,
        "has_significant_changes": bool
    }
    """
    from src.generators.feedback.modification_detector import ModificationDetector

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        detector = ModificationDetector()
        changes = detector.detect_changes(original, modified)
        summary = detector.get_changes_summary()
        has_significant = detector.has_significant_changes()

        # 导出变化到文件
        changes_path = ws.path / "feedback" / "detected_changes.json"
        changes_path.parent.mkdir(parents=True, exist_ok=True)
        detector.export_changes(str(changes_path))

        return json.dumps({
            "changes": changes,
            "summary": summary,
            "has_significant_changes": has_significant,
            "changes_file": str(changes_path)
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def learn_feedback(changes_json: str, slide_index: int, workspace_path: str = "") -> str:
    """从用户修改中学习反馈模式。

    Input: changes_json - 变化JSON, slide_index - 幻灯片索引
    Output: JSON {
        "feedback_summary": dict,
        "recommendations": dict
    }
    """
    from src.generators.feedback.feedback_learner import FeedbackLearner

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        learner = FeedbackLearner(str(ws.path))
        changes = json.loads(changes_json)

        learner.learn_from_changes(changes, slide_index)
        learner.save_feedback()

        summary = learner.get_feedback_summary()
        report = learner.generate_feedback_report()

        # 保存报告
        report_path = ws.path / "feedback" / "feedback_report.md"
        report_path.parent.mkdir(parents=True, exist_ok=True)
        with open(report_path, "w", encoding="utf-8") as f:
            f.write(report)

        return json.dumps({
            "feedback_summary": summary,
            "report_path": str(report_path)
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def get_feedback_patterns(workspace_path: str = "") -> str:
    """获取反馈模式。

    Input: workspace_path - 工作空间路径
    Output: JSON {
        "patterns": dict,
        "history": dict
    }
    """
    from src.generators.feedback.feedback_learner import FeedbackLearner

    ws = Workspace(workspace_path or config.default_workspace)
    ws.ensure_exists()

    try:
        learner = FeedbackLearner(str(ws.path))
        learner.load_feedback()

        summary = learner.get_feedback_summary()

        return json.dumps({
            "patterns": learner.feedback_patterns,
            "summary": summary
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def start_preview_server(port: int = 8765) -> str:
    """启动预览WebSocket服务器，实现HTML页面与Claude Code的实时通信。

    Input: port - 服务器端口号（默认8765）
    Output: JSON {
        "status": str,
        "port": int,
        "message": str
    }
    """
    from src.generators.preview.websocket_server import PreviewServer
    import asyncio

    try:
        server = PreviewServer(port=port)

        # 在后台启动服务器
        async def run_server():
            await server.start()

        # 创建后台任务
        asyncio.create_task(run_server())

        return json.dumps({
            "status": "started",
            "port": port,
            "message": f"WebSocket服务器已启动在端口 {port}。HTML页面将自动连接。"
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


@mcp.tool()
def get_latest_feedback(workspace_path: str = "") -> str:
    """获取最新的用户反馈（从WebSocket服务器）。

    Input: workspace_path - 工作空间路径
    Output: JSON {
        "feedback": list,
        "total": int
    }
    """
    ws = Workspace(workspace_path or config.default_workspace)
    feedback_file = ws.path / "feedback" / "feedback_log.json"

    if not feedback_file.exists():
        return json.dumps({
            "feedback": [],
            "total": 0,
            "message": "暂无反馈"
        }, ensure_ascii=False, indent=2)

    try:
        with open(feedback_file, "r", encoding="utf-8") as f:
            feedback_log = json.load(f)

        return json.dumps({
            "feedback": feedback_log[-10:],  # 最近10条
            "total": len(feedback_log)
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False, indent=2)


if __name__ == "__main__":
    mcp.run()
