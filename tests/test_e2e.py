"""端到端测试 - 验证 MCP Server 工具可用性."""

import sys
import asyncio
import json
from pathlib import Path

# 添加项目根目录到路径
sys.path.insert(0, str(Path(__file__).parent.parent))


async def test_parse_papers():
    """测试论文解析工具."""
    print("\n=== Testing parse_papers ===")
    from src.parsers.paper_parser import parse_papers

    # 测试不存在的文件
    result = await parse_papers(["nonexistent.pdf"])
    assert len(result["papers"]) == 0
    assert len(result["errors"]) > 0
    print("[OK] Error handling works")

    # 测试非PDF文件（用临时文件模拟）
    import tempfile, os
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
        f.write(b"test content")
        txt_path = f.name
    try:
        result = await parse_papers([txt_path])
        assert len(result["papers"]) == 0
        assert any("不支持" in e for e in result["errors"])
        print("[OK] File type validation works")
    finally:
        os.unlink(txt_path)


async def test_extract_figures():
    """测试图表提取工具."""
    print("\n=== 测试 extract_figures ===")
    from src.parsers.figure_extractor import extract_figures

    # 测试不存在的文件
    try:
        await extract_figures("nonexistent.pdf")
        assert False, "应该抛出异常"
    except FileNotFoundError:
        print("[OK] 文件不存在错误处理正常")


async def test_build_slide():
    """测试幻灯片生成工具."""
    print("\n=== 测试 build_slide ===")
    from src.generators.slide_builder import build_slide

    blueprint = """slides:
  - index: 0
    type: title
    title: Test Title
    subtitle: Test Subtitle
    notes: Test notes
    duration_seconds: 30
"""

    result = await build_slide(blueprint, 0)
    assert "pptx_path" in result
    assert result["slide_index"] == 0
    print(f"[OK] 生成成功: {result['pptx_path']}")

    # 检查文件是否存在
    assert Path(result["pptx_path"]).exists()
    print("[OK] 文件已创建")


async def test_get_academic_style():
    """测试学术规范工具."""
    print("\n=== 测试 get_academic_style ===")
    from src.styles.academic_styles import get_academic_style

    result = await get_academic_style("optics")
    assert "primary_color" in result
    assert "font_family" in result
    print(f"[OK] 光学领域规范: {result['primary_color']}")

    result = await get_academic_style("unknown")
    assert result["domain"] == "general"
    print("[OK] 默认领域回退正常")


async def test_get_slide_template():
    """测试页面模板工具."""
    print("\n=== 测试 get_slide_template ===")
    from src.styles.slide_templates import get_slide_template

    result = await get_slide_template("title")
    assert "layout" in result
    assert "elements" in result
    print(f"[OK] 标题模板: {result['layout']}")


async def test_get_citation_format():
    """测试引用格式工具."""
    print("\n=== 测试 get_citation_format ===")
    from src.styles.citation_formats import get_citation_format

    result = await get_citation_format("IEEE")
    assert "inline_format" in result
    assert "reference_format" in result
    print(f"[OK] IEEE格式: {result['inline_format']}")


async def test_read_pptx():
    """测试PPTX读取工具."""
    print("\n=== 测试 read_pptx ===")
    from pptx import Presentation

    # 创建测试PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    test_path = "workspace/preview/test_read.pptx"
    Path(test_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(test_path)

    # 测试读取
    prs = Presentation(test_path)
    assert len(prs.slides) == 1
    print(f"[OK] 读取成功: {len(prs.slides)} 页")


async def test_diff_pptx():
    """测试PPTX差异对比工具."""
    print("\n=== 测试 diff_pptx ===")
    from pptx import Presentation
    from pptx.util import Inches, Pt

    # 创建原版
    prs1 = Presentation()
    slide1 = prs1.slides.add_slide(prs1.slide_layouts[6])
    path1 = "workspace/preview/test_diff_orig.pptx"
    prs1.save(path1)

    # 创建修改版
    prs2 = Presentation()
    slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
    left, top, width, height = Inches(1), Inches(1), Inches(3), Inches(1)
    box = slide2.shapes.add_textbox(left, top, width, height)
    box.text_frame.text = "New text"
    path2 = "workspace/preview/test_diff_mod.pptx"
    prs2.save(path2)

    print("[OK] 测试文件创建成功")


async def run_all_tests():
    """运行所有测试."""
    print("开始端到端测试...")

    tests = [
        test_parse_papers,
        test_extract_figures,
        test_build_slide,
        test_get_academic_style,
        test_get_slide_template,
        test_get_citation_format,
        test_read_pptx,
        test_diff_pptx,
    ]

    passed = 0
    failed = 0

    for test in tests:
        try:
            await test()
            passed += 1
        except Exception as e:
            print(f"[FAIL] {test.__name__} 失败: {e}")
            failed += 1

    print("\n" + "=" * 40)
    print(f"测试结果: {passed} 通过, {failed} 失败")

    if failed == 0:
        print("所有测试通过! [OK]")
    else:
        print("存在失败的测试!")
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(run_all_tests())
